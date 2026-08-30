"""
NEXUS PYQ Training Data Pipeline
=================================
Complete pipeline: Raw Plans → PDF OCR → Question Extraction → CO/Topic Mapping → Training CSV

Workflow:
1. Drop unstructured syllabus/lesson plan docs into data/pyqs/lesson_plans/raw/
2. Drop photographed PYQ PDFs into data/pyqs/raw/
3. Run: python pipeline.py parse-plans  → Parses docs into CO-Module-Topic CSVs
4. Run: python pipeline.py ocr          → OCR all PYQ PDFs → extracted/ folder
5. Run: python pipeline.py extract      → Parse questions from OCR text
6. Run: python pipeline.py map-topics   → Batch Gemini API to classify questions
7. Run: python pipeline.py analyze      → CO frequency analysis
8. Run: python pipeline.py build        → Generate final training CSV
9. Run: python pipeline.py all          → Run entire pipeline end-to-end
"""

import os
import sys
import json
import csv
import re
import time
from pathlib import Path
from collections import defaultdict
import urllib.request

# ============ PATHS ============
DATA_DIR = Path(__file__).parent
BASE_DIR = DATA_DIR.parent.parent

# PYQ Paths
RAW_DIR = DATA_DIR / "raw"              # Drop PYQ PDFs here
EXTRACTED_DIR = DATA_DIR / "extracted"   # OCR output goes here
OUTPUT_DIR = DATA_DIR                    # Final training CSV

# Lesson Plan Paths
LESSON_PLANS_DIR = DATA_DIR / "lesson_plans"
RAW_PLANS_DIR = LESSON_PLANS_DIR / "raw" # Drop unstructured syllabi here

PIPELINE_STATE = DATA_DIR / "pipeline_state.json"

# Ensure dirs exist
for d in [RAW_DIR, EXTRACTED_DIR, LESSON_PLANS_DIR, RAW_PLANS_DIR]:
    d.mkdir(parents=True, exist_ok=True)

# ============ CONFIG ============
GEMINI_MODEL = "gemini-2.5-flash"
GEMINI_MAX_OUTPUT_TOKENS = 8192  # Max tokens per Gemini call
BATCH_SIZE = 15  # Questions per Gemini API call (to stay within limits)


def get_gemini_api_key():
    api_key = os.environ.get("GEMINI_API_KEY", "")
    if not api_key:
        env_path = BASE_DIR / ".env"
        if env_path.exists():
            for line in env_path.read_text(encoding="utf-8").split("\n"):
                if line.startswith("GEMINI_API_KEY"):
                    api_key = line.split("=", 1)[1].strip().strip('"').strip("'")
                    break
    if not api_key:
        print("ERROR: No GEMINI_API_KEY found in environment or .env file")
    return api_key


def load_state():
    if PIPELINE_STATE.exists():
        return json.loads(PIPELINE_STATE.read_text(encoding="utf-8"))
    return {"ocr_done": [], "extracted": [], "mapped": False}


def save_state(state):
    PIPELINE_STATE.write_text(json.dumps(state, indent=2), encoding="utf-8")


# ================================================================
# STEP 0: PARSE PLANS - Unstructured Syllabus to CSV
# ================================================================
def step_parse_plans():
    """Parse raw unstructured lesson plan documents into structured CSVs using Gemini."""
    try:
        import fitz
    except ImportError:
        print("Warning: PyMuPDF not installed, skipping PDF lesson plans.")
        fitz = None
    
    try:
        import docx
    except ImportError:
        print("Warning: python-docx not installed, skipping DOCX lesson plans.")
        docx = None

    raw_files = list(RAW_PLANS_DIR.glob("*.*"))
    valid_exts = {".pdf", ".docx", ".txt"}
    files_to_process = [f for f in raw_files if f.suffix.lower() in valid_exts]

    if not files_to_process:
        print(f"No raw lesson plans found in {RAW_PLANS_DIR}")
        print("Drop your .pdf, .docx, or .txt syllabus documents there.")
        return

    api_key = get_gemini_api_key()
    if not api_key:
        return

    print(f"Found {len(files_to_process)} raw lesson plans to parse...")
    print("=" * 50)
    
    for plan_file in files_to_process:
        print(f"\n  Processing: {plan_file.name}")
        text = ""
        
        if plan_file.suffix.lower() == ".txt":
            text = plan_file.read_text(encoding="utf-8", errors="ignore")
        elif plan_file.suffix.lower() == ".pdf" and fitz:
            doc = fitz.open(plan_file)
            for page in doc:
                text += page.get_text() + "\n"
        elif plan_file.suffix.lower() == ".docx" and docx:
            doc = docx.Document(plan_file)
            text = "\n".join([p.text for p in doc.paragraphs])
        else:
            print(f"    Skipping {plan_file.name} (unsupported or missing library)")
            continue

        if not text.strip():
            print("    Warning: No text extracted.")
            continue

        # Truncate text if too long for API (keep it under ~30k chars)
        text = text[:30000]

        prompt = f"""You are an academic curriculum parser for KIIT University.
        
I am providing a raw, unstructured syllabus or lesson plan document text.
Extract the implicit mapping of Course Outcomes (COs) to Modules, and the Topics covered under each Module.
Pay special attention to tables mapping Sl. No, Topic, and Modules if they exist.

DOCUMENT TEXT:
{text}

OUTPUT FORMAT INSTRUCTIONS:
Return a JSON array of objects representing the curriculum modules. Each object MUST have these EXACT three keys:
- "co": the Course Outcome(s) associated with this module (e.g., "CO1", "CO1, CO2". If not explicitly mapped, try to infer it based on context, else leave as "General").
- "module": the Module name and number (e.g., "Module 1: Foundations of Deep Learning").
- "topics": a comma-separated list of ALL topics covered in this module, combining everything from the lesson plan table (e.g., "Overview of NN, Feed forward network, Gradient descent, Back-propagation algorithm, Multilayer Perceptron").

Respond ONLY with a valid JSON array. No markdown blocks, no other text.
"""
        url = f"https://generativelanguage.googleapis.com/v1beta/models/{GEMINI_MODEL}:generateContent?key={api_key}"
        payload = json.dumps({
            "contents": [{"parts": [{"text": prompt}]}],
            "generationConfig": {
                "temperature": 0.1,
                "responseMimeType": "application/json",
            }
        }).encode("utf-8")
        
        req = urllib.request.Request(url, data=payload, headers={"Content-Type": "application/json"})
        
        try:
            print("    Sending to Gemini for parsing...", end=" ")
            with urllib.request.urlopen(req, timeout=90) as resp:
                data = json.loads(resp.read().decode("utf-8"))
                response_text = data["candidates"][0]["content"]["parts"][0]["text"].strip()
                
                # Strip markdown blocks if present
                if response_text.startswith("```json"):
                    response_text = response_text[7:-3]
                elif response_text.startswith("```"):
                    response_text = response_text[3:-3]
                
                parsed_json = json.loads(response_text)
                
                if not isinstance(parsed_json, list):
                    print("FAILED (Did not return a JSON array)")
                    continue
                
                # Write to CSV
                csv_path = LESSON_PLANS_DIR / f"{plan_file.stem}.csv"
                with open(csv_path, "w", newline="", encoding="utf-8") as f:
                    writer = csv.DictWriter(f, fieldnames=["CO", "Module", "Topics"])
                    writer.writeheader()
                    for item in parsed_json:
                        co = str(item.get("co", "")).replace('"', '""').strip()
                        module = str(item.get("module", "")).replace('"', '""').strip()
                        topics = str(item.get("topics", "")).replace("\n", ", ").replace('"', '""').strip()
                        writer.writerow({"CO": co, "Module": module, "Topics": topics})
                
                print(f"SUCCESS -> Created {csv_path.name} ({len(parsed_json)} modules mapped)")
                
        except Exception as e:
            print(f"FAILED: {e}")

    print(f"\nLesson plan parsing complete! CSVs saved to {LESSON_PLANS_DIR}")

# ================================================================
# STEP 0.5: PARSE PPTs - Professor Notes Extraction
# ================================================================
def step_parse_ppts():
    """Extract text and speaker notes from Class PPTs to use as Knowledge Base."""
    try:
        from pptx import Presentation
    except ImportError:
        print("ERROR: python-pptx not installed. Run: pip install python-pptx")
        return
        
    ppt_dir = DATA_DIR / "knowledge_base" / "raw"
    ppt_dir.mkdir(parents=True, exist_ok=True)
    out_dir = DATA_DIR / "knowledge_base" / "extracted"
    out_dir.mkdir(parents=True, exist_ok=True)
    
    ppts = list(ppt_dir.glob("*.pptx"))
    if not ppts:
        print(f"No professor PPTs found in {ppt_dir}")
        print("Drop .pptx files there to build the context knowledge base.")
        return
        
    print(f"Parsing {len(ppts)} Professor presentations...")
    for ppt_path in ppts:
        try:
            prs = Presentation(ppt_path)
            extracted = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
                notes_text = ""
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    
                full_slide = "\\n--- Slide {i+1} ---\\n" + "\\n".join(slide_text)
                if notes_text.strip():
                    full_slide += "\\n[SPEAKER NOTES]: " + notes_text
                extracted.append(full_slide)
                
            out_file = out_dir / f"{ppt_path.stem}.txt"
            out_file.write_text("\\n\\n".join(extracted), encoding="utf-8")
            print(f"  -> Extracted {lenprs.slides} slides from {ppt_path.name}")
        except Exception as e:
            print(f"  [Error] Failed to process {ppt_path.name}: {e}")

# ================================================================
# STEP 1: OCR - Convert photographed PYQ PDFs to text
# ================================================================
def step_ocr():
    """OCR all PDFs in raw/ folder using PyMuPDF + pytesseract."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        print("ERROR: PyMuPDF not installed. Run: pip install PyMuPDF")
        return

    try:
        import pytesseract
        from PIL import Image
        import io
    except ImportError:
        print("ERROR: pytesseract/Pillow not installed. Run: pip install pytesseract Pillow")
        return

    state = load_state()
    pdfs = list(RAW_DIR.glob("*.pdf")) + list(RAW_DIR.glob("*.PDF"))
    
    if not pdfs:
        print(f"No PDFs found in {RAW_DIR}")
        print(f"Drop your photographed PYQ PDFs there and re-run.")
        return

    print(f"Found {len(pdfs)} PDFs to OCR")
    print("=" * 50)

    for pdf_path in pdfs:
        if pdf_path.name in state["ocr_done"]:
            print(f"  [SKIP] {pdf_path.name} (already done)")
            continue

        print(f"\n  Processing: {pdf_path.name}")
        doc = fitz.open(str(pdf_path))
        full_text = []

        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text().strip()
            
            if len(text) > 50:
                print(f"    Page {page_num + 1}: Direct text ({len(text)} chars)")
                full_text.append(f"--- PAGE {page_num + 1} ---\n{text}")
            else:
                print(f"    Page {page_num + 1}: OCR mode...", end=" ")
                pix = page.get_pixmap(dpi=300)
                img_data = pix.tobytes("png")
                img = Image.open(io.BytesIO(img_data))
                
                ocr_text = pytesseract.image_to_string(img, config='--psm 6')
                print(f"({len(ocr_text)} chars)")
                full_text.append(f"--- PAGE {page_num + 1} (OCR) ---\n{ocr_text}")

        doc.close()

        out_name = pdf_path.stem + ".txt"
        out_path = EXTRACTED_DIR / out_name
        out_path.write_text("\n\n".join(full_text), encoding="utf-8")
        print(f"  Saved: {out_path}")

        state["ocr_done"].append(pdf_path.name)
        save_state(state)

    print(f"\nOCR complete! Extracted texts in: {EXTRACTED_DIR}")


# ================================================================
# STEP 2: EXTRACT - Parse questions from OCR text
# ================================================================
def step_extract():
    """Extract individual questions from OCR text files."""
    txt_files = list(EXTRACTED_DIR.glob("*.txt"))
    if not txt_files:
        print(f"No .txt files in {EXTRACTED_DIR}")
        return

    all_questions = []

    for txt_path in txt_files:
        print(f"\nParsing: {txt_path.name}")
        content = txt_path.read_text(encoding="utf-8")
        
        name_parts = txt_path.stem.replace("-", "_").split("_")
        subject_guess = name_parts[0] if name_parts else "Unknown"
        year_guess = ""
        for part in name_parts:
            if re.match(r"^20\d{2}$", part):
                year_guess = part
                break

        lines = content.split("\n")
        current_question = []
        current_marks = None
        questions_found = []

        for line in lines:
            line = line.strip()
            if not line or line.startswith("---"):
                continue

            q_match = re.match(r'^(?:Q\.?\s*)?(\d+)\s*[.)]\s*(.*)', line, re.IGNORECASE)
            marks_match = re.search(r'\[(\d+)\s*(?:marks?|M)?\]|\((\d+)\s*(?:marks?|M)?\)|(\d+)\s*marks?', line, re.IGNORECASE)

            if q_match:
                if current_question:
                    q_text = " ".join(current_question).strip()
                    if len(q_text) > 10:
                        questions_found.append({
                            "question": q_text,
                            "marks": current_marks or 0,
                            "source_file": txt_path.name,
                        })
                current_question = [q_match.group(2)] if q_match.group(2) else []
                current_marks = None
                if marks_match:
                    current_marks = int(marks_match.group(1) or marks_match.group(2) or marks_match.group(3))
            else:
                if current_question is not None:
                    current_question.append(line)
                if marks_match and current_marks is None:
                    current_marks = int(marks_match.group(1) or marks_match.group(2) or marks_match.group(3))

        if current_question:
            q_text = " ".join(current_question).strip()
            if len(q_text) > 10:
                questions_found.append({
                    "question": q_text,
                    "marks": current_marks or 0,
                    "source_file": txt_path.name,
                })

        for q in questions_found:
            q["subject"] = subject_guess
            q["year"] = year_guess
            q["co"] = ""
            q["topic"] = ""
            q["module"] = ""

        all_questions.extend(questions_found)
        print(f"  Found {len(questions_found)} questions")

    extracted_path = DATA_DIR / "extracted_questions.json"
    with open(extracted_path, "w", encoding="utf-8") as f:
        json.dump(all_questions, f, indent=2, ensure_ascii=False)
    
    print(f"\nExtracted JSON saved to: {extracted_path}")


# ================================================================
# STEP 3: MAP TOPICS - Use Gemini API to classify CO/Topic
# ================================================================
def load_lesson_plans():
    """Load CO → Module → Topic mappings from lesson plan CSVs."""
    plans = {}
    csv_files = list(LESSON_PLANS_DIR.glob("*.csv"))
    
    for csv_path in csv_files:
        subject = csv_path.stem.replace("_", " ").title()
        modules = []
        with open(csv_path, "r", encoding="utf-8") as f:
            reader = csv.DictReader(f)
            for row in reader:
                modules.append({
                    "co": row.get("CO", "").strip(),
                    "module": row.get("Module", "").strip(),
                    "topics": row.get("Topics", "").strip(),
                })
        plans[subject] = modules
        print(f"  Loaded lesson plan: {subject} ({len(modules)} modules)")
    
    return plans


def batch_gemini_classify(questions_batch, lesson_plan_context):
    """Send a batch of questions to Gemini for topic/CO classification."""
    api_key = get_gemini_api_key()
    if not api_key: return None

    questions_text = ""
    for i, q in enumerate(questions_batch):
        questions_text += f"\n[Q{i+1}] ({q.get('marks', '?')} marks) {q['question']}"

    prompt = f"""You are an academic classifier for KIIT University exam questions.

LESSON PLAN / CO MAPPING:
{lesson_plan_context}

QUESTIONS TO CLASSIFY:
{questions_text}

For EACH question, respond with a JSON array. Each element should have:
- "index": the question number (Q1, Q2, etc.)
- "co": the Course Outcome (CO1, CO2, etc.) that best matches
- "topic": a specific topic name from the lesson plan
- "module": the module name/number

Respond ONLY with the JSON array, no other text.
"""
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{GEMINI_MODEL}:generateContent?key={api_key}"
    payload = json.dumps({
        "contents": [{"parts": [{"text": prompt}]}],
        "generationConfig": {"temperature": 0.2, "responseMimeType": "application/json"}
    }).encode("utf-8")

    req = urllib.request.Request(url, data=payload, headers={"Content-Type": "application/json"})
    try:
        with urllib.request.urlopen(req, timeout=60) as resp:
            data = json.loads(resp.read().decode("utf-8"))
            return json.loads(data["candidates"][0]["content"]["parts"][0]["text"])
    except Exception as e:
        print(f"  Gemini API error: {e}")
        return None


def step_map_topics():
    """Map extracted questions to CO/Topics using Gemini API."""
    questions_path = DATA_DIR / "extracted_questions.json"
    if not questions_path.exists():
        print("No extracted_questions.json found. Run extract step first.")
        return

    questions = json.loads(questions_path.read_text(encoding="utf-8"))
    
    print("\nLoading lesson plans...")
    lesson_plans = load_lesson_plans()

    by_subject = defaultdict(list)
    for i, q in enumerate(questions):
        q["_idx"] = i
        by_subject[q.get("subject", "Unknown")].append(q)

    total_api_calls = 0

    for subject, subj_questions in by_subject.items():
        print(f"\n--- {subject} ({len(subj_questions)} questions) ---")

        plan = lesson_plans.get(subject)
        if not plan:
            for plan_subject in lesson_plans:
                if plan_subject.lower() in subject.lower() or subject.lower() in plan_subject.lower():
                    plan = lesson_plans[plan_subject]
                    break

        if plan:
            plan_text = "\n".join([f"  {m['co']}: {m['module']} — Topics: {m['topics']}" for m in plan])
        else:
            plan_text = "(No lesson plan available — classify by general topic)"
            print(f"  WARNING: No lesson plan for '{subject}'. Using general classification.")

        for batch_start in range(0, len(subj_questions), BATCH_SIZE):
            batch = subj_questions[batch_start:batch_start + BATCH_SIZE]
            print(f"  Batch {batch_start+1}-{min(batch_start + BATCH_SIZE, len(subj_questions))}...", end=" ")

            result = batch_gemini_classify(batch, plan_text)
            total_api_calls += 1

            if result and isinstance(result, list):
                for item in result:
                    idx_str = item.get("index", "Q0")
                    idx_num = int(re.search(r"\d+", idx_str).group()) - 1 if re.search(r"\d+", idx_str) else 0
                    if 0 <= idx_num < len(batch):
                        orig_idx = batch[idx_num]["_idx"]
                        questions[orig_idx]["co"] = item.get("co", "")
                        questions[orig_idx]["topic"] = item.get("topic", "")
                        questions[orig_idx]["module"] = item.get("module", "")
                print(f"mapped {len(result)} questions")
            else:
                print("FAILED")
            time.sleep(1)

    for q in questions:
        q.pop("_idx", None)

    mapped_path = DATA_DIR / "mapped_questions.json"
    mapped_path.write_text(json.dumps(questions, indent=2, ensure_ascii=False), encoding="utf-8")
    print(f"\nMapping complete! Saved to {mapped_path}")


# ================================================================
# STEP 4: ANALYZE - CO frequency analysis across years/papers
# ================================================================
def step_analyze():
    mapped_path = DATA_DIR / "mapped_questions.json"
    if mapped_path.exists():
        questions = json.loads(mapped_path.read_text(encoding="utf-8"))
    else:
        csv_path = DATA_DIR / "sample_dataset.csv"
        if not csv_path.exists():
            return
        questions = []
        with open(csv_path, "r", encoding="utf-8") as f:
            for row in csv.DictReader(f):
                questions.append({
                    "subject": row.get("Subject", ""),
                    "question": row.get("Question", ""),
                    "marks": int(row.get("Marks", 0)) if row.get("Marks", "").isdigit() else 0,
                    "year": row.get("Year", ""),
                    "topic": row.get("Topic", ""),
                    "co": row.get("CO", ""),
                })

    print(f"Analyzing {len(questions)} questions...\n")
    
    # Subject Dist
    by_subject = defaultdict(list)
    for q in questions: by_subject[q.get("subject", "Unknown")].append(q)

    # CO Repetition
    co_data = [q for q in questions if q.get("co")]
    if co_data:
        co_matrix = defaultdict(lambda: defaultdict(lambda: defaultdict(int)))
        for q in co_data:
            co_matrix[q.get("subject", "?")][str(q.get("year", "?"))][q.get("co", "?")] += 1
            
        print("\nCO REPETITION PATTERNS:")
        for subj in sorted(co_matrix.keys()):
            co_totals, co_years = defaultdict(int), defaultdict(set)
            for year, cos in co_matrix[subj].items():
                for co, count in cos.items():
                    co_totals[co] += count
                    co_years[co].add(year)
            print(f"\n  {subj}:")
            for co in sorted(co_totals.keys()):
                print(f"    {co}: {co_totals[co]} questions across {len(co_years[co])} years")

    # Topic Dist
    by_topic = defaultdict(int)
    for q in questions:
        if t := q.get("topic"): by_topic[t] += 1
        
    analysis_path = DATA_DIR / "analysis.json"
    analysis_path.write_text(json.dumps({
        "subjects": {s: len(qs) for s, qs in by_subject.items()},
        "topics": dict(by_topic),
    }, indent=2), encoding="utf-8")
    print("\nAnalysis saved.")


# ================================================================
# STEP 5: BUILD - Generate final training CSV
# ================================================================
def step_build():
    for source in ["mapped_questions.json", "extracted_questions.json"]:
        source_path = DATA_DIR / source
        if source_path.exists():
            questions = json.loads(source_path.read_text(encoding="utf-8"))
            break
    else:
        return

    valid = []
    for q in questions:
        question_text = q.get("question", "").strip()
        subject = q.get("subject", "").strip()
        if len(question_text) < 15 or not subject or subject == "Unknown": continue
        marks = int(q.get("marks", 5)) if str(q.get("marks", "")).isdigit() else 5
        valid.append({
            "Subject": subject,
            "Question": question_text,
            "Marks": marks,
            "Year": q.get("year", "2023"),
            "Topic": q.get("topic", "General"),
        })

    existing_path = DATA_DIR / "sample_dataset.csv"
    existing = list(csv.DictReader(open(existing_path, "r", encoding="utf-8"))) if existing_path.exists() else []

    existing_texts = {row["Question"].lower().strip()[:50] for row in existing}
    new_qs = [v for v in valid if v["Question"].lower().strip()[:50] not in existing_texts]

    with open(existing_path, "w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=["Subject", "Question", "Marks", "Year", "Topic"])
        writer.writeheader()
        writer.writerows(existing + new_qs)

    print(f"\nCSV built: {len(existing)} existing + {len(new_qs)} new = {len(existing) + len(new_qs)} total.")


# ================================================================
# MAIN
# ================================================================
def print_usage():
    print("""
NEXUS PYQ Training Data Pipeline
=================================

Usage: python pipeline.py <step>

Steps (run in order):
  parse-plans  Parses raw lesson plans (PDF/DOCX/TXT) to structured CSV using Gemini
  ocr          OCR photographed PYQ PDFs
  extract      Parse questions from OCR text
  map-topics   Batch Gemini API → classify CO/Topic
  analyze      CO frequency analysis
  build        Generate final training CSV
  all          Run entire pipeline
""")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print_usage()
        sys.exit(0)

    cmd = sys.argv[1].lower()

    if cmd == "parse-plans": step_parse_plans()
    elif cmd == "ocr": step_ocr()
    elif cmd == "extract": step_extract()
    elif cmd in ("map-topics", "map"): step_map_topics()
    elif cmd in ("analyze", "analysis"): step_analyze()
    elif cmd == "build": step_build()
    elif cmd == "all":
        print("Running full pipeline...\n")
        step_parse_plans()
        print("\n" + "=" * 60 + "\n")
        step_ocr()
        print("\n" + "=" * 60 + "\n")
        step_extract()
        print("\n" + "=" * 60 + "\n")
        step_map_topics()
        print("\n" + "=" * 60 + "\n")
        step_analyze()
        print("\n" + "=" * 60 + "\n")
        step_build()
    else:
        print(f"Unknown command: {cmd}")
        print_usage()
